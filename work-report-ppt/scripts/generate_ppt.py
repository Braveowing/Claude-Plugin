#!/usr/bin/env python3
"""
Work Report PPT Generator (English Output - Compact Format)
Generates professional work report presentations in English from Markdown content using SCQA framework.
Supports compact format: 1 page for simple reports, 1-2 pages for complex reports, with image support.
"""

import argparse
import re
from pathlib import Path
from typing import List, Dict, Optional
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as RgbColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is required. Install it with: pip install python-pptx")
    exit(1)


class SlideColors:
    """Color scheme for professional presentations."""

    TITLE = RgbColor(23, 43, 77)  # Dark blue-gray
    HEADING = RgbColor(0, 82, 204)  # Corporate blue
    BODY = RgbColor(31, 31, 31)  # Dark gray
    LIGHT = RgbColor(102, 102, 102)  # Light gray
    WHITE = RgbColor(255, 255, 255)
    ACCENT = RgbColor(0, 82, 204)
    SUCCESS = RgbColor(54, 179, 126)
    WARNING = RgbColor(255, 171, 0)
    ERROR = RgbColor(255, 86, 48)


class FontSizes:
    """Font sizes for different text elements (compact format)."""

    TITLE = 44
    HEADING = 32
    SUBHEADING = 24
    BODY = 20
    SMALL = 16


class CompactPPTGenerator:
    """Generator for SCQA-based work report presentations in English - Compact Format."""

    def __init__(self, template_path: Optional[str] = None):
        """Initialize generator with optional template."""
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
            print(f"Using template: {template_path}")
        else:
            self.prs = Presentation()
            # Set slide dimensions to 16:9
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(5.625)
            print("Using default template (16:9 aspect ratio)")

        self.colors = SlideColors()
        self.fonts = FontSizes()

    def add_compact_simple_slide(self, sections: Dict):
        """Add single-slide compact report (Simple Report - 1 Page)."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # Report Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = sections.get('title', 'Work Report')
        self._format_text_frame(title_frame, self.fonts.TITLE, self.colors.TITLE, bold=True, center=True)

        # Presenter and Date (compact)
        date_str = sections.get('date') or datetime.now().strftime("%B %d, %Y")
        presenter = sections.get('presenter', '')
        if presenter:
            meta_text = f"Presenter: {presenter} | Date: {date_str}"
        else:
            meta_text = f"Date: {date_str}"
        meta_box = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(5), Inches(0.4))
        meta_frame = meta_box.text_frame
        meta_frame.text = meta_text
        self._format_text_frame(meta_frame, self.fonts.SMALL, self.colors.LIGHT, center=True)

        # Executive Summary (compact)
        exec_top = Inches(1.6)
        exec_box = slide.shapes.add_textbox(Inches(0.5), exec_top, Inches(9), Inches(0.6))
        exec_frame = exec_box.text_frame
        p = exec_frame.add_paragraph()
        p.text = "Executive Summary"
        self._format_paragraph(p, self.fonts.HEADING, self.colors.HEADING, bold=True)

        # Executive summary bullets (compact, inline)
        bullets = sections.get('executive_summary', [])[:3]
        if bullets:
            bullets_text = " | ".join(f"• {b}" for b in bullets)
            bullet_box = slide.shapes.add_textbox(Inches(0.5), exec_top + Inches(0.5), Inches(9), Inches(0.5))
            bullet_frame = bullet_box.text_frame
            bullet_frame.text = bullets_text
            self._format_text_frame(bullet_frame, self.fonts.SUBHEADING, self.colors.BODY)


        # Problem & Root Cause (compact)
        problem_top = exec_top + Inches(1.4)
        heading_box = slide.shapes.add_textbox(Inches(0.5), problem_top, Inches(9), Inches(0.4))
        heading_frame = heading_box.text_frame
        heading_frame.text = "Problem & Root Cause"
        self._format_text_frame(heading_frame, self.fonts.HEADING, self.colors.HEADING, bold=True)

        # Problem description
        problem = " | ".join(sections.get('complication', sections.get('situation', [])))[:200]
        root_cause = " | ".join(sections.get('root_cause', []))[:100]
        problem_box = slide.shapes.add_textbox(Inches(0.5), problem_top + Inches(0.5), Inches(9), Inches(0.6))
        problem_frame = problem_box.text_frame
        problem_frame.word_wrap = True
        p = problem_frame.add_paragraph()
        p.text = f"Problem: {problem[:100]}"
        self._format_paragraph(p, self.fonts.SUBHEADING, self.colors.BODY)
        p = problem_frame.add_paragraph()
        p.text = f"Root Cause: {root_cause[:80]}"
        self._format_paragraph(p, self.fonts.SUBHEADING, self.colors.BODY)

        # Solution Options (compact table)
        options = sections.get('solution_options', [])[:3]
        if options:
            self._add_compact_solution_table(slide, options, problem_top + Inches(1.6))

        # Recommendation & Decision (compact)
        rec_top = problem_top + Inches(3)
        rec_box = slide.shapes.add_textbox(Inches(0.5), rec_top, Inches(9), Inches(0.4))
        rec_frame = rec_box.text_frame
        rec_frame.text = "Recommendation"
        self._format_text_frame(rec_frame, self.fonts.HEADING, self.colors.ACCENT, bold=True)

        rec_text = sections.get('recommendation', '')
        just_text = " | ".join(sections.get('justification', []))[:100]
        dec_text = sections.get('decision_required', '')

        rec_detail_box = slide.shapes.add_textbox(Inches(0.5), rec_top + Inches(0.5), Inches(9), Inches(1))
        rec_detail_frame = rec_detail_box.text_frame
        rec_detail_frame.word_wrap = True
        p = rec_detail_frame.add_paragraph()
        p.text = f"Recommended: {rec_text[:80]}"
        self._format_paragraph(p, self.fonts.SUBHEADING, self.colors.ACCENT, bold=True)
        p = rec_detail_frame.add_paragraph()
        p.text = f"Justification: {just_text[:80]}"
        self._format_paragraph(p, self.fonts.SUBHEADING, self.colors.BODY)
        if dec_text:
            p = rec_detail_frame.add_paragraph()
            p.text = f"Decision: {dec_text[:80]}"
            self._format_paragraph(p, self.fonts.SUBHEADING, self.colors.ERROR, bold=True)

    def _add_compact_solution_table(self, slide, options: List[Dict], top: float):
        """Add compact solution comparison table."""
        num_options = len(options)
        col_width = (Inches(9)) / num_options

        for i, option in enumerate(options):
            left = Inches(0.5) + (col_width * i)

            # Option name
            name_box = slide.shapes.add_textbox(left, top, col_width - Inches(0.1), Inches(0.4))
            name_frame = name_box.text_frame
            name_frame.text = f"Option {i+1}: {option.get('name', 'Unnamed')[:30]}"
            self._format_text_frame(name_frame, self.fonts.SUBHEADING, self.colors.HEADING, bold=True)

            # Pros (compact)
            pros = " • ".join(option.get('pros', [])[:2])
            pros_box = slide.shapes.add_textbox(left, top + Inches(0.5), col_width - Inches(0.15), Inches(0.5))
            pros_frame = pros_box.text_frame
            pros_frame.text = f"Pros: {pros[:60]}"
            self._format_text_frame(pros_frame, self.fonts.BODY, self.colors.BODY)

            # Cons (compact)
            cons = " • ".join(option.get('cons', [])[:1])
            cons_box = slide.shapes.add_textbox(left, top + Inches(1.1), col_width - Inches(0.15), Inches(0.4))
            cons_frame = cons_box.text_frame
            cons_frame.text = f"Cons: {cons[:40]}"
            self._format_text_frame(cons_frame, self.fonts.BODY, self.colors.BODY)

            # Cost & Timeline (single line)
            cost = option.get('cost', 'N/A')
            timeline = option.get('timeline', 'N/A')
            info_box = slide.shapes.add_textbox(left, top + Inches(1.6), col_width - Inches(0.15), Inches(0.5))
            info_frame = info_box.text_frame
            info_frame.text = f"Cost: {cost[:15]} | Time: {timeline[:15]}"
            self._format_text_frame(info_frame, self.fonts.SUBHEADING, self.colors.LIGHT)

    def add_compact_complex_slides(self, sections: Dict, images: Dict):
        """Add 2-slide compact report (Complex Report - 2 Pages)."""

        # SLIDE 1: Problem Analysis & Solutions
        blank_layout = self.prs.slide_layouts[6]
        slide1 = self.prs.slides.add_slide(blank_layout)

        # Report Title
        title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = sections.get('title', 'Work Report')
        self._format_text_frame(title_frame, self.fonts.TITLE, self.colors.TITLE, bold=True, center=True)

        # Presenter and Date (compact)
        date_str = sections.get('date') or datetime.now().strftime("%B %d, %Y")
        presenter = sections.get('presenter', '')
        if presenter:
            meta_text = f"Presenter: {presenter} | Date: {date_str}"
        else:
            meta_text = f"Date: {date_str}"
        meta_box = slide1.shapes.add_textbox(Inches(5), Inches(0.9), Inches(5), Inches(0.35))
        meta_frame = meta_box.text_frame
        meta_frame.text = meta_text
        self._format_text_frame(meta_frame, self.fonts.SMALL, self.colors.LIGHT, center=True)

        # Executive Summary
        exec_top = Inches(1.5)
        exec_box = slide1.shapes.add_textbox(Inches(0.5), exec_top, Inches(9), Inches(0.5))
        exec_frame = exec_box.text_frame
        p = exec_frame.add_paragraph()
        p.text = "Executive Summary"
        self._format_paragraph(p, self.fonts.HEADING, self.colors.HEADING, bold=True)

        bullets = sections.get('executive_summary', [])[:4]
        bullets_text = "\n".join(f"• {b}" for b in bullets)
        bullet_box = slide1.shapes.add_textbox(Inches(0.5), exec_top + Inches(0.5), Inches(4.25), Inches(0.7))
        bullet_frame = bullet_box.text_frame
        bullet_frame.text = bullets_text
        self._format_text_frame(bullet_frame, self.fonts.BODY, self.colors.BODY)

        # Right side - Problem Analysis
        problem_top = exec_top + Inches(1.4)
        heading_box = slide1.shapes.add_textbox(Inches(0.5), problem_top, Inches(4.25), Inches(0.35))
        heading_frame = heading_box.text_frame
        heading_frame.text = "Problem Analysis"
        self._format_text_frame(heading_frame, self.fonts.HEADING, self.colors.HEADING, bold=True)

        # Situation (compact)
        situation = " | ".join(sections.get('situation', []))[:60]
        sit_box = slide1.shapes.add_textbox(Inches(0.5), problem_top + Inches(0.45), Inches(4.25), Inches(0.4))
        sit_frame = sit_box.text_frame
        sit_frame.text = f"Situation: {situation}"
        self._format_text_frame(sit_frame, self.fonts.SUBHEADING, self.colors.BODY)

        # Problem (compact)
        problem = " | ".join(sections.get('complication', []))[:60]
        prob_box = slide1.shapes.add_textbox(Inches(0.5), problem_top + Inches(0.9), Inches(4.25), Inches(0.4))
        prob_frame = prob_box.text_frame
        prob_frame.text = f"Problem: {problem}"
        self._format_text_frame(prob_frame, self.fonts.SUBHEADING, self.colors.BODY)

        # Root Cause (compact)
        root_cause = "\n".join(sections.get('root_cause', [])[:3])
        rc_box = slide1.shapes.add_textbox(Inches(0.5), problem_top + Inches(1.35), Inches(4.25), Inches(0.7))
        rc_frame = rc_box.text_frame
        rc_frame.text = f"Root Cause:\n{root_cause}"
        self._format_text_frame(rc_frame, self.fonts.SUBHEADING, self.colors.BODY)

        # Solution Options (right side)
        solutions_top = problem_top
        options = sections.get('solution_options', [])[:3]
        self._add_compact_solution_table_two_col(slide1, options, solutions_top + Inches(2.4))

        # SLIDE 2: Recommendation & Decision
        slide2 = self.prs.slides.add_slide(blank_layout)

        # Recommendation Title
        rec_title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        rec_title_frame = rec_title_box.text_frame
        rec_title_frame.text = "Recommendation & Decision"
        self._format_text_frame(rec_title_frame, self.fonts.TITLE, self.colors.ACCENT, bold=True, center=True)

        # Divider
        self._add_divider(slide2, Inches(0.5), Inches(0.8), Inches(9))

        # Recommended
        rec_top = Inches(1.1)
        rec_text = sections.get('recommendation', '')

        rec_box = slide2.shapes.add_textbox(Inches(0.5), rec_top, Inches(9), Inches(0.4))
        rec_frame = rec_box.text_frame
        rec_frame.text = f"Recommended: {rec_text[:100]}"
        self._format_text_frame(rec_frame, self.fonts.SUBHEADING, self.colors.ACCENT, bold=True)

        # Justification
        just_top = rec_top + Inches(0.6)
        just_items = sections.get('justification', [])[:3]
        just_box = slide2.shapes.add_textbox(Inches(0.5), just_top, Inches(4.25), Inches(0.8))
        just_frame = just_box.text_frame
        just_frame.text = "Justification:"
        self._format_paragraph(just_frame.paragraphs[0], self.fonts.SUBHEADING, self.colors.HEADING, bold=True)
        for just in just_items:
            p = just_frame.add_paragraph()
            p.text = f"• {just[:80]}"
            self._format_paragraph(p, self.fonts.BODY, self.colors.BODY)

        # Risk Mitigation
        risk_top = just_top
        risk_items = sections.get('risk_mitigation', [])[:3]
        risk_box = slide2.shapes.add_textbox(Inches(5.5), risk_top, Inches(4.25), Inches(0.8))
        risk_frame = risk_box.text_frame
        risk_frame.text = "Risk Mitigation:"
        self._format_paragraph(risk_frame.paragraphs[0], self.fonts.SUBHEADING, self.colors.WARNING, bold=True)
        for risk in risk_items:
            p = risk_frame.add_paragraph()
            p.text = f"• {risk[:80]}"
            self._format_paragraph(p, self.fonts.BODY, self.colors.BODY)

        # Implementation Plan
        impl_top = just_top + Inches(1.3)
        impl_items = sections.get('next_steps', [])[:4]
        impl_box = slide2.shapes.add_textbox(Inches(0.5), impl_top, Inches(4.25), Inches(0.8))
        impl_frame = impl_box.text_frame
        impl_frame.text = "Implementation Plan:"
        self._format_paragraph(impl_frame.paragraphs[0], self.fonts.SUBHEADING, self.colors.HEADING, bold=True)
        for impl in impl_items:
            p = impl_frame.add_paragraph()
            p.text = f"• {impl[:80]}"
            self._format_paragraph(p, self.fonts.BODY, self.colors.BODY)

        # Decision Required
        dec_top = impl_top
        dec_text = sections.get('decision_required', '')

        dec_box = slide2.shapes.add_textbox(Inches(5.5), dec_top, Inches(4.25), Inches(0.8))
        dec_frame = dec_box.text_frame
        dec_frame.text = "Decision Required:"
        self._format_paragraph(dec_frame.paragraphs[0], self.fonts.SUBHEADING, self.colors.ERROR, bold=True)

        p = dec_frame.add_paragraph()
        p.text = f"• {dec_text[:150]}"
        self._format_paragraph(p, self.fonts.BODY, self.colors.BODY)

    def _add_compact_solution_table_two_col(self, slide, options: List[Dict], top: float):
        """Add 2-column compact solution comparison."""
        # Create 2 columns for 3 options (top:2, bottom:1)
        col_width = Inches(2.1)

        # Top row - Options 1 & 2
        for i in range(min(2, len(options))):
            opt = options[i]
            left = Inches(0.5) + (col_width * i)

            # Option name
            name_box = slide.shapes.add_textbox(left, top, col_width - Inches(0.1), Inches(0.35))
            name_frame = name_box.text_frame
            name_frame.text = f"Option {i+1}: {opt.get('name', 'Unnamed')[:25]}"
            self._format_text_frame(name_frame, self.fonts.SUBHEADING, self.colors.HEADING, bold=True)

            # Pros
            pros = " • ".join(opt.get('pros', [])[:2])
            pros_box = slide.shapes.add_textbox(left, top + Inches(0.4), col_width - Inches(0.1), Inches(0.45))
            pros_frame = pros_box.text_frame
            pros_frame.text = f"Pros:\n{pros[:50]}"
            self._format_text_frame(pros_frame, self.fonts.BODY, self.colors.BODY)

            # Cons & Info
            cons = opt.get('cons', [''])[:1][0] if opt.get('cons') else 'N/A'
            cost = opt.get('cost', 'N/A')
            timeline = opt.get('timeline', 'N/A')
            info_box = slide.shapes.add_textbox(left, top + Inches(0.9), col_width - Inches(0.1), Inches(0.6))
            info_frame = info_box.text_frame
            info_frame.text = f"Cons: {cons[:30]}\nCost: {cost}\nTime: {timeline}"
            self._format_text_frame(info_frame, self.fonts.BODY, self.colors.BODY)

        # Bottom row - Option 3 (if exists)
        if len(options) > 2:
            opt = options[2]
            left = Inches(0.5) + col_width

            # Center it under options 1&2
            name_box = slide.shapes.add_textbox(left, top + Inches(1.9), col_width * 2 - Inches(0.1), Inches(0.35))
            name_frame = name_box.text_frame
            name_frame.text = f"Option 3: {opt.get('name', 'Unnamed')[:25]}"
            self._format_text_frame(name_frame, self.fonts.SUBHEADING, self.colors.HEADING, bold=True)

            # Pros
            pros = " • ".join(opt.get('pros', [])[:2])
            pros_box = slide.shapes.add_textbox(left, top + Inches(2.3), col_width * 2 - Inches(0.1), Inches(0.45))
            pros_frame = pros_box.text_frame
            pros_frame.text = f"Pros:\n{pros[:50]}"
            self._format_text_frame(pros_frame, self.fonts.BODY, self.colors.BODY)

            # Cons & Info
            cons = opt.get('cons', [''])[:1][0] if opt.get('cons') else 'N/A'
            cost = opt.get('cost', 'N/A')
            timeline = opt.get('timeline', 'N/A')
            info_box = slide.shapes.add_textbox(left, top + Inches(2.8), col_width * 2 - Inches(0.1), Inches(0.6))
            info_frame = info_box.text_frame
            info_frame.text = f"Cons: {cons[:30]}\nCost: {cost}\nTime: {timeline}"
            self._format_text_frame(info_frame, self.fonts.BODY, self.colors.BODY)

    def add_image_to_slide(self, slide, image_path: str, left: float, top: float, width: float, height: float):
        """Add image to slide."""
        try:
            if Path(image_path).exists():
                slide.shapes.add_picture(image_path, left, top, width, height)
                print(f"Added image: {image_path}")
            else:
                print(f"Warning: Image not found: {image_path}")
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")

    def _add_divider(self, slide, left: float, top: float, width: float):
        """Add horizontal divider line."""
        shape = slide.shapes.add_shape(MSO_SHAPE.LINE, left, top)
        line = shape.line
        line.color.rgb = self.colors.LIGHT
        line.width = Pt(1)
        shape.width = width
        shape.height = Pt(1)

    def _format_text_frame(self, text_frame, font_size: int, color, bold: bool = False, center: bool = False):
        """Format a text frame's paragraphs."""
        for paragraph in text_frame.paragraphs:
            self._format_paragraph(paragraph, font_size, color, bold, center)

    def _format_paragraph(self, paragraph, font_size: int, color, bold: bool = False, center: bool = False):
        """Format a single paragraph."""
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT

    def save(self, output_path: str):
        """Save presentation to file."""
        self.prs.save(output_path)
        print(f"Presentation saved to: {output_path}")


def parse_markdown(content: str) -> Dict:
    """Parse markdown content and extract SCQA sections with solution options and images."""
    sections = {
        'title': '',
        'executive_summary': [],
        'situation': [],
        'complication': [],
        'root_cause': [],
        'question': '',
        'solution_options': [],
        'recommendation': '',
        'justification': [],
        'risk_mitigation': [],
        'next_steps': [],
        'decision_required': '',
        'presenter': '',
        'date': '',
        'report_type': 'simple',
        'images': {}
    }

    current_section = None
    current_option = {}
    lines = content.split('\n')

    for line in lines:
        line = line.strip()

        # Detect section headers
        if line.startswith('# '):
            sections['title'] = line[2:].strip()
        elif line.startswith('## Executive Summary') or line.lower().startswith('## executive'):
            current_section = 'executive_summary'
        elif line.startswith('## Situation') or line.lower().startswith('## 现状'):
            current_section = 'situation'
        elif line.startswith('## Complication') or line.lower().startswith('## 冲突'):
            current_section = 'complication'
        elif line.startswith('## Root Cause') or line.lower().startswith('## 根因'):
            current_section = 'root_cause'
        elif line.startswith('## Key Question') or line.lower().startswith('## 问题'):
            current_section = 'question'
        elif line.startswith('## Solution Options') or line.lower().startswith('## 解决方案'):
            current_section = 'solution_options'
        elif line.startswith('## Recommendation') or line.lower().startswith('## 推荐'):
            current_section = 'recommendation'
        elif line.startswith('## Next Steps') or line.lower().startswith('## 下一步'):
            current_section = 'next_steps'
        elif line.startswith('## Decision Required') or line.lower().startswith('## 决策'):
            current_section = 'decision_required'
        elif line.startswith('## Images') or line.lower().startswith('## 图片'):
            current_section = 'images'
        elif line.startswith('### Option'):
            # Save previous option if exists
            if current_option:
                sections['solution_options'].append(current_option)
            # Start new option
            option_name = line.replace('### Option', '').replace('###', '').strip()
            current_option = {'name': option_name, 'pros': [], 'cons': []}
            current_section = 'option_details'
        elif line.startswith('- '):
            bullet = line[2:].strip()
            if current_section in ['executive_summary', 'situation', 'complication',
                               'root_cause', 'justification', 'risk_mitigation', 'next_steps']:
                sections[current_section].append(bullet)
            elif current_section == 'option_details' and current_option:
                if bullet.lower().startswith('pros:'):
                    continue
                else:
                    current_option['pros'].append(bullet)
            elif current_section == 'option_cons' and current_option:
                current_option['cons'].append(bullet)
        elif line.lower().startswith('pros:'):
            if current_section == 'option_details':
                continue
        elif line.lower().startswith('cons:'):
            if current_section == 'option_details':
                current_section = 'option_cons'
        elif ':' in line and current_section == 'option_details':
            # Parse option details like Cost:, Timeline:, Description:
            key, value = line.split(':', 1)
            key = key.strip().lower()
            value = value.strip()
            if key in ['cost', 'timeline', 'description']:
                current_option[key] = value
        elif line.startswith('Image') and current_section == 'images':
            # Parse image: Image1: path/to/image.png
            if ':' in line:
                key, path = line.split(':', 1)
                sections['images'][key.strip()] = path.strip()
        elif line.startswith('Presenter:'):
            sections['presenter'] = line.split(':', 1)[1].strip()
        elif line.startswith('Date:'):
            sections['date'] = line.split(':', 1)[1].strip()
        elif line.startswith('Report Type:'):
            sections['report_type'] = line.split(':', 1)[1].strip().lower()
        elif line and not line.startswith('#'):
            # Non-bullet content
            if current_section == 'question':
                sections['question'] = line
            elif current_section == 'recommendation':
                if line.lower().startswith('recommended:'):
                    sections['recommendation'] = line.split(':', 1)[1].strip()
                elif line.lower().startswith('justification:'):
                    current_section = 'justification'
                elif line.lower().startswith('risk mitigation:'):
                    current_section = 'risk_mitigation'
                else:
                    sections['recommendation'] = line
            elif current_section in ['decision_required']:
                sections[current_section] += line + " "

    # Save last option if exists
    if current_option:
        sections['solution_options'].append(current_option)

    return sections


def generate_from_markdown(input_path: str, output_path: str, template_path: Optional[str] = None):
    """Generate PPT from markdown file."""
    # Read markdown content
    with open(input_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Parse content
    sections = parse_markdown(content)

    # Create generator
    generator = CompactPPTGenerator(template_path)

    # Generate based on report type
    report_type = sections.get('report_type', 'simple').lower()

    if report_type == 'simple':
        print("Generating simple report (1 page)...")
        generator.add_compact_simple_slide(sections)
    else:
        print("Generating complex report (2 pages)...")
        generator.add_compact_complex_slides(sections, sections.get('images', {}))

    # Save
    generator.save(output_path)
    print(f"Generated {len(generator.prs.slides)} slide(s)")


def main():
    parser = argparse.ArgumentParser(
        description='Generate work report PPT in English from Markdown using SCQA framework (Compact Format)'
    )
    parser.add_argument('--input', '-i', required=True,
                        help='Input Markdown file path')
    parser.add_argument('--output', '-o', required=True,
                        help='Output PPTX file path')
    parser.add_argument('--template', '-t',
                        help='Template PPTX file path')

    args = parser.parse_args()

    # Validate input file exists
    if not Path(args.input).exists():
        print(f"Error: Input file not found: {args.input}")
        exit(1)

    # Generate PPT
    print(f"Generating presentation from: {args.input}")
    generate_from_markdown(args.input, args.output, args.template)
    print("Done!")


if __name__ == '__main__':
    main()
