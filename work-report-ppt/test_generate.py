import argparse
import sys
from pathlib import Path

# Test import
print(f"Python: {sys.executable}")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RgbColor
    from pptx.enum.shapes import MSO_SHAPE
    print("pptx import: SUCCESS")
except ImportError as e:
    print(f"pptx import: FAILED - {e}")
    sys.exit(1)

# Read and parse markdown
input_path = "D:/AI Project/Claude code/skills/work-report-ppt/content.md"
output_path = "D:/AI Project/Claude code/skills/work-report-ppt/MPA_Feature_Delivery_Report.pptx"

with open(input_path, 'r', encoding='utf-8') as f:
    content = f.read()

# Simple parsing for testing
sections = {
    'title': 'MPA Feature Delivery Decision Report',
    'date': '2026-03-01',
    'executive_summary': ['MPA delayed due to collision issues', 'Decision: waiver vs. postpone to R7', 'Safety concerns with waiver option'],
    'complication': ['Perception model collides with barriers, carts, metal objects'],
    'root_cause': ['MMT unwilling to optimize R6 architecture', 'R7 fix available in Px96 (2026 Q3)'],
    'solution_options': [
        {
            'name': 'Waiver Acceptance',
            'pros': ['On-time delivery', 'Meets timeline'],
            'cons': ['Safety risks', 'Customer complaints'],
            'cost': 'Minimal',
            'timeline': 'Px82.25 (2026 Q1)'
        },
        {
            'name': 'Postpone to R7',
            'pros': ['All safety resolved', 'Better UX', 'No complaint risk'],
            'cons': ['6-month delay', 'Misses milestone'],
            'cost': 'Delay cost',
            'timeline': 'Px96 (2026 Q3)'
        }
    ],
    'recommendation': 'Option 2 - Postpone to R7',
    'justification': ['Safety is non-negotiable', 'R7 provides permanent solution', 'Better to delay than deliver unsafe'],
    'decision_required': 'Approve postpone to Px96 (R7 release)'
}

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add compact simple slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = sections['title']
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.color.rgb = RgbColor(23, 43, 77)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Date
date_str = sections['date']
meta_box = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(5), Inches(0.4))
meta_frame = meta_box.text_frame
meta_frame.text = f"Date: {date_str}"
meta_frame.paragraphs[0].font.size = Pt(16)
meta_frame.paragraphs[0].font.color.rgb = RgbColor(102, 102, 102)
meta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Executive Summary
exec_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(0.6))
exec_frame = exec_box.text_frame
exec_frame.text = "Executive Summary"
exec_frame.paragraphs[0].font.size = Pt(32)
exec_frame.paragraphs[0].font.color.rgb = RgbColor(0, 82, 204)
exec_frame.paragraphs[0].font.bold = True

bullets = " | ".join(f"• {b}" for b in sections['executive_summary'])
bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(9), Inches(0.5))
bullet_frame = bullet_box.text_frame
bullet_frame.text = bullets
bullet_frame.paragraphs[0].font.size = Pt(20)
bullet_frame.paragraphs[0].font.color.rgb = RgbColor(31, 31, 31)

# Problem
prob_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.6))
prob_frame = prob_box.text_frame
prob_frame.text = "Problem & Root Cause"
prob_frame.paragraphs[0].font.size = Pt(32)
prob_frame.paragraphs[0].font.color.rgb = RgbColor(0, 82, 204)
prob_frame.paragraphs[0].font.bold = True

prob_detail_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(9), Inches(0.6))
prob_detail_frame = prob_detail_box.text_frame
prob_detail_frame.word_wrap = True
p = prob_detail_frame.add_paragraph()
p.text = f"Problem: {sections['complication'][0]}"
p.font.size = Pt(20)
p.font.color.rgb = RgbColor(31, 31, 31)
p = prob_detail_frame.add_paragraph()
p.text = f"Root Cause: {sections['root_cause'][0]}"
p.font.size = Pt(20)
p.font.color.rgb = RgbColor(31, 31, 31)

# Recommendation
rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.4))
rec_frame = rec_box.text_frame
rec_frame.text = "Recommendation"
rec_frame.paragraphs[0].font.size = Pt(32)
rec_frame.paragraphs[0].font.color.rgb = RgbColor(0, 82, 204)
rec_frame.paragraphs[0].font.bold = True

rec_text = f"Recommended: {sections['recommendation']}"
just_text = f"Justification: {sections['justification'][0]}"
dec_text = f"Decision: {sections['decision_required']}"

rec_detail_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.9))
rec_detail_frame = rec_detail_box.text_frame
rec_detail_frame.word_wrap = True
p = rec_detail_frame.add_paragraph()
p.text = rec_text
p.font.size = Pt(20)
p.font.color.rgb = RgbColor(0, 82, 204)
p.font.bold = True
p = rec_detail_frame.add_paragraph()
p.text = just_text
p.font.size = Pt(20)
p.font.color.rgb = RgbColor(31, 31, 31)
p = rec_detail_frame.add_paragraph()
p.text = dec_text
p.font.size = Pt(20)
p.font.color.rgb = RgbColor(255, 86, 48)
p.font.bold = True

# Save
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Generated {len(prs.slides)} slide(s)")
