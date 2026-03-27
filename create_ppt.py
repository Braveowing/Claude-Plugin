from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
import io

# MB Corporate Colors
MB_BLACK = RGBColor(0, 0, 0)
MB_WHITE = RGBColor(255, 255, 255)
MB_SILVER = RGBColor(192, 192, 192)
MB_DARK_GRAY = RGBColor(64, 64, 64)
MB_LIGHT_GRAY = RGBColor(245, 245, 245)
MB_BLUE = RGBColor(0, 100, 180)  # Subtle MB blue

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(blank_layout)

# ===== HEADER SECTION =====
# Add header background
header_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(13.333), Inches(1.0)
)
header_bg.fill.solid()
header_bg.fill.fore_color.rgb = RGBColor(30, 30, 30)
header_bg.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.25),
    Inches(12.333), Inches(0.5)
)
title_frame = title_box.text_frame
title_frame.text = "CIVIC FUP2 Change Request: road_ownership Signal Integration"
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = MB_WHITE
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Subtitle
subtitle_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.65),
    Inches(12.333), Inches(0.3)
)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Critical Dependency for P2P Functionality | IDC R16 / 2026-3.0 OTA"
subtitle_frame.paragraphs[0].font.size = Pt(14)
subtitle_frame.paragraphs[0].font.color.rgb = MB_SILVER
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# ===== CONTEXT BOX =====
context_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.3), Inches(1.2),
    Inches(12.733), Inches(0.65)
)
context_box.fill.solid()
context_box.fill.fore_color.rgb = RGBColor(230, 240, 250)
context_box.line.color.rgb = MB_BLUE
context_box.line.width = Pt(1.5)

context_text = slide.shapes.add_textbox(
    Inches(0.5), Inches(1.3),
    Inches(12.333), Inches(0.5)
)
tf = context_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Context: 2026-3.0 OTA has no planned release → CR delayed to 2026-4.0 OTA → P2P functionality delayed"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(40, 60, 80)
p.alignment = PP_ALIGN.LEFT

# ===== SECTION 1: WHY CHANGE (LEFT SIDE) =====
section1_title = slide.shapes.add_textbox(
    Inches(0.3), Inches(2.0),
    Inches(6.3), Inches(0.3)
)
tf = section1_title.text_frame
tf.text = "Why This Change Is Necessary"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = MB_DARK_GRAY

# Functional Impact Box
func_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.3), Inches(2.4),
    Inches(6.3), Inches(2.2)
)
func_box.fill.solid()
func_box.fill.fore_color.rgb = MB_LIGHT_GRAY
func_box.line.color.rgb = RGBColor(180, 180, 180)
func_box.line.width = Pt(1)

# Functional Impact Title
func_title = slide.shapes.add_textbox(
    Inches(0.5), Inches(2.5),
    Inches(6.1), Inches(0.3)
)
tf = func_title.text_frame
tf.text = "Functional Impact"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = MB_BLACK

# Create parking lot diagram
img = Image.new('RGB', (600, 200), color='white')
draw = ImageDraw.Draw(img)

# Draw parking lot boundary
draw.rectangle([50, 30, 550, 170], outline='gray', width=2)

# Draw entrance/exit points (misaligned)
draw.line([(100, 30), (100, 170)], fill='red', width=2)  # Left boundary
draw.rectangle([95, 40, 105, 80], fill='red')  # Actual entrance
draw.text((70, 15), "Actual Entry", fill='red', font=ImageFont.load_default())

draw.rectangle([295, 120, 305, 160], fill='orange')  # Expected exit (wrong position)
draw.text((260, 175), "Expected Exit (Wrong)", fill='orange', font=ImageFont.load_default())

draw.rectangle([445, 40, 455, 80], fill='green')  # Correct exit position
draw.text((420, 15), "Correct Exit", fill='green', font=ImageFont.load_default())

# Draw arrow showing misalignment
draw.line([(300, 140), (450, 60)], fill='blue', width=2)
draw.polygon([(450, 60), (440, 75), (460, 75)], fill='blue')

# Draw SD route (wrong)
draw.arc([120, 50, 280, 150], start=0, end=180, fill='red', width=2)
draw.text((150, 85), "SD Route", fill='red', font=ImageFont.load_default())

# Draw MPA route
draw.arc([320, 50, 440, 130], start=0, end=180, fill='green', width=2)
draw.text((340, 85), "MPA Route", fill='green', font=ImageFont.load_default())

# Label
draw.text((200, 180), "Parking Lot Area", fill='black', font=ImageFont.load_default())
draw.text((50, 5), "❌ Without road_ownership: Entry/Exit Misalignment", fill='red', font=ImageFont.load_default())

# Save to bytes
img_bytes = io.BytesIO()
img.save(img_bytes, format='PNG')
img_bytes.seek(0)

# Add to slide
slide.shapes.add_picture(img_bytes, Inches(0.4), Inches(2.85), width=Inches(6.1))

# Functional Impact Text
func_text = slide.shapes.add_textbox(
    Inches(0.5), Inches(4.7),
    Inches(6.1), Inches(0.8)
)
tf = func_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• Entry/exit positioning deviation\n• Reduced SD-MPA route stitching success rate\n• Lower L2PP ↔ MPA switching reliability\n• Significant user experience degradation"
p.font.size = Pt(11)
p.font.color.rgb = MB_DARK_GRAY
p.level = 0

# ===== SECTION 2: SIGNAL LINK COMPARISON (RIGHT SIDE) =====
section2_title = slide.shapes.add_textbox(
    Inches(6.8), Inches(2.0),
    Inches(6.3), Inches(0.3)
)
tf = section2_title.text_frame
tf.text = "Signal Delivery Path Analysis"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = MB_DARK_GRAY

# Cloud Link Box (NOT SUITABLE)
cloud_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.8), Inches(2.4),
    Inches(6.2), Inches(1.1)
)
cloud_box.fill.solid()
cloud_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
cloud_box.line.color.rgb = RGBColor(200, 50, 50)
cloud_box.line.width = Pt(2)

cloud_title = slide.shapes.add_textbox(
    Inches(7.0), Inches(2.5),
    Inches(6.0), Inches(0.3)
)
tf = cloud_title.text_frame
tf.text = "❌ Cloud SD Route Path (NOT SUITABLE)"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(180, 0, 0)

cloud_text = slide.shapes.add_textbox(
    Inches(7.0), Inches(2.85),
    Inches(6.0), Inches(0.6)
)
tf = cloud_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• SD route only available during active navigation\n• Users may enter parking lot WITHOUT navigation active\n• No real-time road_ownership in non-navigation scenarios"
p.font.size = Pt(10)
p.font.color.rgb = MB_DARK_GRAY

# Vehicle Bus Link Box (RECOMMENDED)
bus_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.8), Inches(3.65),
    Inches(6.2), Inches(1.1)
)
bus_box.fill.solid()
bus_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
bus_box.line.color.rgb = RGBColor(50, 150, 50)
bus_box.line.width = Pt(2)

bus_title = slide.shapes.add_textbox(
    Inches(7.0), Inches(3.75),
    Inches(6.0), Inches(0.3)
)
tf = bus_title.text_frame
tf.text = "✓ Vehicle Bus Path (RECOMMENDED)"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0, 120, 0)

bus_text = slide.shapes.add_textbox(
    Inches(7.0), Inches(4.1),
    Inches(6.0), Inches(0.6)
)
tf = bus_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• MX NDS compiles Amap SDK axf link data\n• Real-time road_ownership via existing bus signal\n• Works in ALL scenarios (navigation & non-navigation)"
p.font.size = Pt(10)
p.font.color.rgb = MB_DARK_GRAY

# Additional reason box
reason_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.8), Inches(4.9),
    Inches(6.2), Inches(0.75)
)
reason_box.fill.solid()
reason_box.fill.fore_color.rgb = RGBColor(255, 250, 230)
reason_box.line.color.rgb = RGBColor(200, 150, 50)
reason_box.line.width = Pt(1)

reason_text = slide.shapes.add_textbox(
    Inches(7.0), Inches(5.0),
    Inches(6.0), Inches(0.6)
)
tf = reason_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Key Insight: Other existing navigation signals cannot characterize parking lot area information"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = RGBColor(150, 100, 0)

# ===== SOLUTION SECTION =====
solution_title = slide.shapes.add_textbox(
    Inches(0.3), Inches(5.6),
    Inches(5.0), Inches(0.3)
)
tf = solution_title.text_frame
tf.text = "Solution & Implementation"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = MB_DARK_GRAY

solution_text = slide.shapes.add_textbox(
    Inches(0.3), Inches(5.95),
    Inches(6.3), Inches(0.6)
)
tf = solution_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• MX NDS development (based on Amap SDK compilation)\n• IDC interface layer adaptation\n• Reuse existing bus signal (no new signal required)\n• Effort: Low | Risk: Controllable"
p.font.size = Pt(11)
p.font.color.rgb = MB_DARK_GRAY

# ===== TIMING BOX =====
timing_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.3), Inches(6.65),
    Inches(6.3), Inches(0.7)
)
timing_box.fill.solid()
timing_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
timing_box.line.color.rgb = MB_SILVER
timing_box.line.width = Pt(1)

timing_text = slide.shapes.add_textbox(
    Inches(0.5), Inches(6.75),
    Inches(6.1), Inches(0.5)
)
tf = timing_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Timing: Raised now because MMT completed R7 architecture change in March 2026, which finalized this solution and dependencies"
p.font.size = Pt(10)
p.font.italic = True
p.font.color.rgb = RGBColor(80, 80, 80)

# ===== RECOMMENDATION BOX (HIGHLIGHTED) =====
rec_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.8), Inches(5.8),
    Inches(6.2), Inches(1.0)
)
rec_box.fill.solid()
rec_box.fill.fore_color.rgb = MB_BLUE
rec_box.line.fill.background()

rec_title = slide.shapes.add_textbox(
    Inches(7.0), Inches(5.9),
    Inches(6.0), Inches(0.3)
)
tf = rec_title.text_frame
tf.text = "RECOMMENDATION"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = MB_WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

rec_text = slide.shapes.add_textbox(
    Inches(7.0), Inches(6.25),
    Inches(6.0), Inches(0.5)
)
tf = rec_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Add CIVIC FUP2 version to R16 release\nEnable road_ownership signal integration for P2P"
p.font.size = Pt(12)
p.font.color.rgb = MB_WHITE
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = r"D:\AI\CIVIC_FUP2_CR_Presentation.pptx"
prs.save(output_path)
print(f"PPT created successfully: {output_path}")
