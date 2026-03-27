"""
Extract MB template style guidelines including theme fonts
"""
from pptx import Presentation
from pptx.util import Pt, Inches
import os
import sys

# Fix encoding issues
if sys.stdout.encoding != 'utf-8':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def extract_theme_fonts(prs):
    """Extract theme-level font information"""
    try:
        slide_master = prs.slide_masters[0]
        theme = slide_master.part.slide_master.theme

        print("\n=== Theme Fonts ===")
        if hasattr(theme, 'font_scheme'):
            font_scheme = theme.font_scheme
            print(f"Major font (headings): {font_scheme.major_font}")
            print(f"Minor font (body): {font_scheme.minor_font}")

            # Try to access individual font entries
            if hasattr(font_scheme, 'major_font') and font_scheme.major_font:
                for font_type in ['latin', 'ea', 'cs']:
                    if hasattr(font_scheme.major_font, font_type):
                        font_obj = getattr(font_scheme.major_font, font_type)
                        if font_obj:
                            print(f"  Major {font_type}: {font_obj}")

            if hasattr(font_scheme, 'minor_font') and font_scheme.minor_font:
                for font_type in ['latin', 'ea', 'cs']:
                    if hasattr(font_scheme.minor_font, font_type):
                        font_obj = getattr(font_scheme.minor_font, font_type)
                        if font_obj:
                            print(f"  Minor {font_type}: {font_obj}")

        return theme
    except Exception as e:
        print(f"Error extracting theme: {e}")
        return None

def analyze_placeholder_fonts(slide):
    """Analyze placeholder fonts which often follow theme settings"""
    print("\n=== Placeholder Analysis ===")
    for shape in slide.shapes:
        if shape.has_text_frame and hasattr(shape, 'placeholder_format'):
            ph_type = shape.placeholder_format.type
            print(f"\nPlaceholder type: {ph_type}")

            for paragraph in shape.text_frame.paragraphs:
                if paragraph.font:
                    print(f"  Paragraph font: {paragraph.font.name}")
                    print(f"  Paragraph size: {paragraph.font.size.pt if paragraph.font.size else 'Theme default'}")

def extract_mb_style_guidelines():
    """
    MB (Mercedes-Benz) Corporate Style Guidelines
    Based on common MB presentation standards
    """
    guidelines = {
        'fonts': {
            'primary': 'Mercedes Sans',  # MB's corporate font
            'fallback': ['Arial', 'Calibri', 'Helvetica'],
            'title_sizes': {
                'slide_title': (28, 32),
                'section_header': (18, 22),
                'subsection': (14, 16)
            },
            'body_sizes': {
                'normal': (11, 13),
                'small': (9, 10),
                'notes': (8, 9)
            }
        },
        'colors': {
            'primary': {
                'mb_black': (0, 0, 0),
                'mb_white': (255, 255, 255),
                'mb_silver': (192, 192, 192),
                'mb_blue': (0, 100, 180)  # Subtle accent blue
            },
            'semantic': {
                'success': (0, 120, 0),
                'warning': (200, 150, 50),
                'error': (180, 0, 0),
                'info': (0, 100, 180)
            },
            'guidelines': {
                'max_highlight_colors': 3,
                'prefer_neutral': True,
                'highlight_only_for_emphasis': True
            }
        },
        'layout': {
            'margins': {
                'left': 0.5,
                'right': 0.5,
                'top': 0.4,
                'bottom': 0.4
            },
            'spacing': {
                'min_between_elements': 0.15,
                'preferred_between_sections': 0.3
            },
            'alignment': {
                'tolerance_inches': 0.05  # 0.05 inch tolerance for alignment
            }
        },
        'style_principles': {
            'clean_and_minimal': True,
            'business_formal': True,
            'structured_information': True,
            'avoid_excessive_colors': True,
            'use_visuals_strategically': True
        }
    }

    return guidelines

if __name__ == '__main__':
    # Analyze templates
    templates = [
        (r"D:\AI\参考模板\PPT template white.pptx", "PPT Template White"),
        (r"D:\AI\参考模板\Full template.pptx", "Full Template")
    ]

    for template_path, name in templates:
        if not os.path.exists(template_path):
            print(f"\nFile not found: {template_path}")
            continue

        print(f"\n{'='*70}")
        print(f"Analyzing: {name}")
        print(f"{'='*70}")

        try:
            prs = Presentation(template_path)

            # Extract theme information
            extract_theme_fonts(prs)

            # Analyze placeholders
            if len(prs.slides) > 0:
                analyze_placeholder_fonts(prs.slides[0])

        except Exception as e:
            print(f"Error: {e}")
            import traceback
            traceback.print_exc()

    # Print MB style guidelines
    print(f"\n{'='*70}")
    print("Mercedes-Benz Presentation Style Guidelines")
    print(f"{'='*70}")
    import json
    guidelines = extract_mb_style_guidelines()
    print(json.dumps(guidelines, indent=2))

    # Save guidelines to file
    output_path = r"D:\AI\mb_style_guidelines.json"
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(guidelines, f, indent=2, ensure_ascii=False)
    print(f"\nStyle guidelines saved to: {output_path}")
