"""
Analyze MB template structure to understand font usage
"""
from pptx import Presentation
from pptx.util import Pt, Inches
import os

def analyze_template(template_path):
    """Analyze template structure"""
    if not os.path.exists(template_path):
        print(f"File not found: {template_path}")
        return

    print(f"\n{'='*60}")
    print(f"Analyzing: {os.path.basename(template_path)}")
    print(f"{'='*60}")

    prs = Presentation(template_path)
    print(f"Slide width: {prs.slide_width.inches} inches")
    print(f"Slide height: {prs.slide_height.inches} inches")
    print(f"Number of slides: {len(prs.slides)}")

    for idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {idx + 1} ---")
        print(f"Number of shapes: {len(slide.shapes)}")

        for shape_idx, shape in enumerate(slide.shapes):
            print(f"\nShape {shape_idx + 1}:")
            print(f"  Type: {shape.shape_type}")
            print(f"  Position: ({shape.left.inches:.2f}, {shape.top.inches:.2f}) inches")
            print(f"  Size: {shape.width.inches:.2f} x {shape.height.inches:.2f} inches")

            if shape.has_text_frame:
                print(f"  Has text frame: Yes")
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    print(f"    Paragraph {para_idx + 1}:")
                    if paragraph.text:
                        print(f"      Text: '{paragraph.text[:50]}...'")

                    for run_idx, run in enumerate(paragraph.runs):
                        if run.text:
                            print(f"        Run {run_idx + 1}: '{run.text[:30]}'")
                            if run.font:
                                print(f"          Font name: {run.font.name}")
                                print(f"          Font size: {run.font.size.pt if run.font.size else 'None'}")
                                print(f"          Bold: {run.font.bold}")
                                print(f"          Italic: {run.font.italic}")

if __name__ == '__main__':
    templates = [
        r"D:\AI\参考模板\PPT template white.pptx",
        r"D:\AI\参考模板\Full template.pptx"
    ]

    for template in templates:
        try:
            analyze_template(template)
        except Exception as e:
            print(f"Error analyzing {template}: {e}")
            import traceback
            traceback.print_exc()
