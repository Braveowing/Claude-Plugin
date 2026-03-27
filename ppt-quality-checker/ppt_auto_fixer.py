"""
PPT Layout Fixer - Automatically fix layout and style issues
根据检查结果自动修复PPT布局和样式问题
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import shutil
from datetime import datetime
from typing import List, Dict, Tuple

class PPTLayoutFixer:
    """Automatically fix common PPT layout and style issues"""

    def __init__(self, mb_guidelines_path: str = None):
        """Initialize fixer with MB guidelines"""
        self.mb_guidelines = self._load_guidelines(mb_guidelines_path)
        self.fixes_applied = []
        self.fixes_failed = []

    def _load_guidelines(self, path: str) -> Dict:
        """Load MB style guidelines"""
        if path and os.path.exists(path):
            import json
            with open(path, 'r', encoding='utf-8') as f:
                return json.load(f)
        else:
            # Default MB guidelines
            return {
                'fonts': {
                    'primary': 'Mercedes Sans',
                    'fallback': ['Arial', 'Calibri', 'Helvetica']
                },
                'colors': {
                    'primary': {
                        'mb_black': [0, 0, 0],
                        'mb_white': [255, 255, 255],
                        'mb_silver': [192, 192, 192],
                        'mb_blue': [0, 100, 180]
                    }
                },
                'layout': {
                    'margins': {
                        'left': 0.5,
                        'right': 0.5,
                        'top': 0.4,
                        'bottom': 0.4
                    },
                    'spacing': {
                        'min_between_elements': 0.15
                    }
                }
            }

    def backup_presentation(self, ppt_path: str) -> str:
        """Create backup before fixing"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup_path = ppt_path.replace('.pptx', f'_backup_{timestamp}.pptx')
        shutil.copy2(ppt_path, backup_path)
        return backup_path

    def fix_presentation(self, ppt_path: str, fix_types: List[str] = None,
                        create_backup: bool = True,
                        save_to_new_file: bool = False,
                        new_filename_suffix: str = "_fixed") -> Tuple[bool, List[str]]:
        """
        Fix presentation issues

        Args:
            ppt_path: Path to PPT file
            fix_types: Types of fixes to apply (None = all)
                - 'margins': Fix margin violations
                - 'fonts': Fix non-MB fonts
                - 'colors': Reduce color usage
                - 'spacing': Fix spacing issues
                - 'alignment': Fix alignment issues
            create_backup: Whether to create backup before fixing
            save_to_new_file: If True, saves to new file instead of overwriting
            new_filename_suffix: Suffix for new filename (e.g., "_fixed")

        Returns:
            (success, list_of_fixes)
        """
        # Load presentation
        try:
            prs = Presentation(ppt_path)
        except Exception as e:
            return False, [f"Failed to load PPT: {e}"]

        # Create backup
        backup_path = None
        if create_backup:
                backup_path = self.backup_presentation(ppt_path)
                self.fixes_applied.append(f"Backup created: {os.path.basename(backup_path)}")

        # Default to all fixes
        if fix_types is None:
            fix_types = ['margins', 'fonts', 'colors', 'spacing', 'alignment']

        # Apply fixes
        try:
            if 'margins' in fix_types:
                self._fix_margins(prs)

            if 'fonts' in fix_types:
                self._fix_fonts(prs)

            if 'colors' in fix_types:
                self._fix_colors(prs)

            if 'spacing' in fix_types:
                self._fix_spacing(prs)

            if 'alignment' in fix_types:
                self._fix_alignment(prs)

            # Save fixed presentation
            if save_to_new_file:
                # Save to new file
                from pathlib import Path
                from datetime import datetime

                ppt_path_obj = Path(ppt_path)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                new_filename = f"{ppt_path_obj.stem}{new_filename_suffix}_{timestamp}{ppt_path_obj.suffix}"
                save_path = str(ppt_path_obj.parent / new_filename)
            else:
                save_path = ppt_path

            prs.save(save_path)
            self.fixes_applied.append(f"✅ Presentation saved: {save_path}")

            return True, self.fixes_applied

        except Exception as e:
            self.fixes_failed.append(f"Fix failed: {e}")
            return False, self.fixes_failed

    def _fix_margins(self, prs: Presentation):
        """Fix margin violations"""
        margins = self.mb_guidelines['layout']['margins']
        left_min = Inches(margins['left'])
        right_min = Inches(margins['right'])
        top_min = Inches(margins['top'])
        bottom_min = Inches(margins['bottom'])

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if not hasattr(shape, 'left'):
                    continue

                fixed = False

                # Fix left margin
                if shape.left < left_min:
                    shape.left = left_min
                    fixed = True

                # Fix top margin
                if shape.top < top_min:
                    shape.top = top_min
                    fixed = True

                # Fix right margin
                if shape.left + shape.width > slide_width - right_min:
                    # Reduce width if possible, otherwise move left
                    max_width = slide_width - right_min - shape.left
                    if max_width > Inches(1):  # Reasonable minimum width
                        shape.width = max_width
                    else:
                        shape.left = max(Inches(0), slide_width - right_min - shape.width)
                    fixed = True

                # Fix bottom margin
                if shape.top + shape.height > slide_height - bottom_min:
                    # Reduce height if possible, otherwise move up
                    max_height = slide_height - bottom_min - shape.top
                    if max_height > Inches(0.5):  # Reasonable minimum height
                        shape.height = max_height
                    else:
                        shape.top = max(Inches(0), slide_height - bottom_min - shape.height)
                    fixed = True

                if fixed:
                    self.fixes_applied.append(f"Slide {slide_idx} - Fixed margins for {shape.name}")

    def _fix_fonts(self, prs: Presentation):
        """Fix non-MB fonts"""
        approved_fonts = [self.mb_guidelines['fonts']['primary']] + \
                         self.mb_guidelines['fonts']['fallback']

        # Primary fallback font (usually Arial)
        fallback_font = self.mb_guidelines['fonts']['fallback'][0] if \
                       self.mb_guidelines['fonts']['fallback'] else 'Arial'

        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font and run.font.name:
                            if run.font.name not in approved_fonts:
                                old_font = run.font.name
                                run.font.name = fallback_font
                                self.fixes_applied.append(
                                    f"Slide {slide_idx} - Changed font from '{old_font}' to '{fallback_font}' in {shape.name}"
                                )

    def _fix_colors(self, prs: Presentation):
        """
        Fix color usage - reduce to max 3 accent colors
        This is a conservative fix that doesn't change core design
        """
        # Get MB approved colors
        mb_colors = []
        for color_name, rgb in self.mb_guidelines['colors']['primary'].items():
            mb_colors.append(RGBColor(*rgb))

        # For each slide, count colors and suggest reduction
        # Note: This is complex to auto-fix without breaking design
        # So we'll add to information rather than automatically change

        for slide_idx, slide in enumerate(prs.slides, 1):
            color_count = {}

            for shape in slide.shapes:
                # Check fill color
                if hasattr(shape, 'fill') and shape.fill:
                    try:
                        if shape.fill.type is not None:
                            fill_color = shape.fill.fore_color.rgb
                            if fill_color:
                                color_key = (fill_color[0], fill_color[1], fill_color[2])
                                color_count[color_key] = color_count.get(color_key, 0) + 1
                    except:
                        pass

            # If too many colors, add to report (don't auto-fix to preserve design)
            if len(color_count) > 3:
                self.fixes_applied.append(
                    f"Slide {slide_idx} - ⚠️  Uses {len(color_count)} colors (recommended max: 3). "
                    "Manual review suggested for color reduction."
                )

    def _fix_spacing(self, prs: Presentation):
        """Fix spacing issues between elements"""
        min_spacing = Inches(self.mb_guidelines['layout']['spacing']['min_between_elements'])

        for slide_idx, slide in enumerate(prs.slides, 1):
            shapes_list = list(slide.shapes)

            # Sort shapes by position (top-to-bottom, left-to-right)
            shapes_list.sort(key=lambda s: (s.top if hasattr(s, 'top') else 0,
                                           s.left if hasattr(s, 'left') else 0))

            for i in range(len(shapes_list) - 1):
                shape1 = shapes_list[i]
                shape2 = shapes_list[i + 1]

                if not hasattr(shape1, 'top') or not hasattr(shape2, 'top'):
                    continue

                # Check vertical spacing
                if abs(shape1.left - shape2.left) < Inches(1):  # Roughly aligned vertically
                    spacing = shape2.top - (shape1.top + shape1.height)

                    if 0 < spacing < min_spacing:
                        # Move shape2 down to create proper spacing
                        shape2.top = shape1.top + shape1.height + min_spacing
                        self.fixes_applied.append(
                            f"Slide {slide_idx} - Increased vertical spacing between {shape1.name} and {shape2.name}"
                        )

    def _fix_alignment(self, prs: Presentation):
        """Fix alignment issues"""
        tolerance = Inches(0.05)  # 0.05 inch tolerance

        for slide_idx, slide in enumerate(prs.slides, 1):
            shapes_list = [s for s in slide.shapes if hasattr(s, 'left')]

            # Group shapes by approximate vertical position
            position_groups = {}
            for shape in shapes_list:
                pos_key = round(shape.top / Inches(0.5))  # Group by 0.5 inch bands
                if pos_key not in position_groups:
                    position_groups[pos_key] = []
                position_groups[pos_key].append(shape)

            # Align shapes in each group
            for pos_key, group in position_groups.items():
                if len(group) < 2:
                    continue

                # Find most common left edge (assumed to be the intended alignment)
                left_edges = [s.left for s in group]
                target_left = max(set(left_edges), key=left_edges.count)

                # Align all shapes to target
                for shape in group:
                    if abs(shape.left - target_left) > tolerance:
                        old_left = shape.left
                        shape.left = target_left
                        self.fixes_applied.append(
                            f"Slide {slide_idx} - Aligned {shape.name} (moved {abs(old_left - target_left).inches:.2f}in)"
                        )

    def generate_fix_report(self) -> str:
        """Generate a report of fixes applied"""
        report = []
        report.append("="*80)
        report.append("PPT FIX REPORT")
        report.append("="*80)
        report.append(f"Total fixes applied: {len(self.fixes_applied)}")
        report.append(f"Total fixes failed: {len(self.fixes_failed)}")
        report.append("")

        if self.fixes_applied:
            report.append("✅ FIXES APPLIED:")
            for i, fix in enumerate(self.fixes_applied, 1):
                report.append(f"  {i}. {fix}")
            report.append("")

        if self.fixes_failed:
            report.append("❌ FIXES FAILED:")
            for i, fix in enumerate(self.fixes_failed, 1):
                report.append(f"  {i}. {fix}")
            report.append("")

        report.append("="*80)

        return "\n".join(report)


def main():
    """Command-line usage"""
    import sys
    import io

    # Fix encoding for Windows console
    if sys.stdout.encoding != 'utf-8':
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    if len(sys.argv) < 2:
        print("Usage: python ppt_auto_fixer.py <ppt_file_path> [fix_types] [--overwrite]")
        print("\nFix types (comma-separated):")
        print("  - margins: Fix margin violations")
        print("  - fonts: Fix non-MB fonts")
        print("  - colors: Reduce color usage")
        print("  - spacing: Fix spacing issues")
        print("  - alignment: Fix alignment issues")
        print("  - all: Apply all fixes (default)")
        print("\nOptions:")
        print("  --overwrite    Save to original file (default: save to new file)")
        sys.exit(1)

    ppt_path = sys.argv[1]
    overwrite = '--overwrite' in sys.argv

    # Parse fix types
    fix_types = None
    for arg in sys.argv[2:]:
        if arg != '--overwrite':
            fix_types_arg = arg.lower()
            if fix_types_arg == 'all':
                fix_types = None  # All fixes
            else:
                fix_types = [f.strip() for f in fix_types_arg.split(',')]
            break

    # Run fixer
    print(f"\n🔧 Fixing: {os.path.basename(ppt_path)}")
    if overwrite:
        print("Mode: Overwrite original file (with backup)")
    else:
        print("Mode: Save to new file")

    fixer = PPTLayoutFixer(r"D:\AI\mb_style_guidelines.json")
    success, fixes = fixer.fix_presentation(
        ppt_path,
        fix_types,
        create_backup=True,
        save_to_new_file=not overwrite
    )

    # Generate report
    report = fixer.generate_fix_report()
    print("\n" + report)

    if success:
        print(f"\n✅ Successfully applied {len(fixer.fixes_applied)} fixes!")
    else:
        print(f"\n❌ Fix failed: {fixes[0] if fixes else 'Unknown error'}")
        sys.exit(1)


if __name__ == '__main__':
    main()
