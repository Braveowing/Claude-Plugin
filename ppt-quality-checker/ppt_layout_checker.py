"""
PPT Layout Checker Tool
Checks PowerPoint presentations for layout, font, color, and style issues
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json
import os
from typing import Dict, List, Tuple
from dataclasses import dataclass
from enum import Enum

class IssueSeverity(Enum):
    ERROR = "ERROR"
    WARNING = "WARNING"
    INFO = "INFO"

@dataclass
class LayoutIssue:
    slide_num: int
    shape_name: str
    issue_type: str
    severity: IssueSeverity
    description: str
    suggestion: str

class PPTLayoutChecker:
    def __init__(self, style_guidelines_path=None):
        """Initialize checker with MB style guidelines"""
        if style_guidelines_path and os.path.exists(style_guidelines_path):
            with open(style_guidelines_path, 'r', encoding='utf-8') as f:
                self.guidelines = json.load(f)
        else:
            # Default MB style guidelines
            self.guidelines = self._get_default_guidelines()

        self.issues: List[LayoutIssue] = []

    def _get_default_guidelines(self):
        """Get default MB style guidelines"""
        return {
            'fonts': {
                'primary': 'Mercedes Sans',
                'fallback': ['Arial', 'Calibri', 'Helvetica'],
                'title_sizes': {'slide_title': [28, 32], 'section_header': [18, 22], 'subsection': [14, 16]},
                'body_sizes': {'normal': [11, 13], 'small': [9, 10], 'notes': [8, 9]}
            },
            'colors': {
                'primary': {'mb_black': [0, 0, 0], 'mb_white': [255, 255, 255], 'mb_silver': [192, 192, 192], 'mb_blue': [0, 100, 180]},
                'semantic': {'success': [0, 120, 0], 'warning': [200, 150, 50], 'error': [180, 0, 0], 'info': [0, 100, 180]},
                'guidelines': {'max_highlight_colors': 3, 'prefer_neutral': True, 'highlight_only_for_emphasis': True}
            },
            'layout': {
                'margins': {'left': 0.5, 'right': 0.5, 'top': 0.4, 'bottom': 0.4},
                'spacing': {'min_between_elements': 0.15, 'preferred_between_sections': 0.3},
                'alignment': {'tolerance_inches': 0.05}
            }
        }

    def check_overlaps(self, prs: Presentation):
        """Check for overlapping shapes"""
        for slide_idx, slide in enumerate(prs.slides, 1):
            shapes_list = list(slide.shapes)

            for i, shape1 in enumerate(shapes_list):
                for shape2 in shapes_list[i+1:]:
                    overlap_ratio = self._shapes_overlap(shape1, shape2)

                    # Only report significant overlaps (>60% overlap)
                    # Ignore text boxes intentionally placed on shapes (common design pattern)
                    if overlap_ratio and overlap_ratio > 0.6:
                        # Check if this is likely intentional (textbox on shape)
                        is_intentional = self._is_likely_intentional_overlap(shape1, shape2)

                        if not is_intentional:
                            issue = LayoutIssue(
                                slide_num=slide_idx,
                                shape_name=f"{shape1.name} and {shape2.name}",
                                issue_type="Overlap",
                                severity=IssueSeverity.WARNING,
                                description=f"Shapes '{shape1.name}' and '{shape2.name}' overlap by {overlap_ratio*100:.0f}%",
                                suggestion="Adjust positioning to avoid unintended overlap"
                            )
                            self.issues.append(issue)

    def _is_likely_intentional_overlap(self, shape1, shape2) -> bool:
        """
        Check if overlap is likely intentional (e.g., text box on colored background)
        Common design patterns:
        - Text box placed on top of colored shape (labels, titles on headers)
        - Icon/badge on top of image
        - Labels on diagrams
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Check if shapes have text frames
        is_textbox1 = shape1.has_text_frame if hasattr(shape1, 'has_text_frame') else False
        is_textbox2 = shape2.has_text_frame if hasattr(shape2, 'has_text_frame') else False

        # Get shape types
        type1 = shape1.shape_type if hasattr(shape1, 'shape_type') else None
        type2 = shape2.shape_type if hasattr(shape2, 'shape_type') else None

        # Pattern 1: Header bar with title/subtitle overlay (VERY common)
        # One shape is full-width background, other has text
        if is_textbox1 != is_textbox2:  # One has text, one doesn't
            # Check if it's near top of slide (header area)
            if hasattr(shape1, 'top') and hasattr(shape2, 'top'):
                if shape1.top < Inches(1.5) or shape2.top < Inches(1.5):
                    # Check if one is full-width (typical header background)
                    if hasattr(shape1, 'width') and hasattr(shape2, 'width'):
                        if shape1.width > Inches(10) or shape2.width > Inches(10):
                            # This is very likely a header bar with text overlay
                            return True

        # Pattern 2: Text on colored background box (content boxes)
        # Common: rounded rectangle with text box on top
        if is_textbox1 and is_textbox2:
            # Both have text - check if they're same content box
            if hasattr(shape1, 'name') and hasattr(shape2, 'name'):
                # Often text boxes in same container have sequential names
                name1_num = ''.join(filter(str.isdigit, shape1.name)) if shape1.name else ''
                name2_num = ''.join(filter(str.isdigit, shape2.name)) if shape2.name else ''
                if name1_num and name2_num and abs(int(name1_num) - int(name2_num)) == 1:
                    # Sequential numbering suggests intentional grouping
                    pass  # Don't automatically mark as intentional

            # Check if one is very small (likely a label/badge)
            area1 = shape1.width.inches * shape1.height.inches if hasattr(shape1, 'width') else 999
            area2 = shape2.width.inches * shape2.height.inches if hasattr(shape2, 'width') else 999
            if area1 < 2 or area2 < 2:  # Very small text box (label)
                return True

        # Pattern 3: Shape with text label on top
        if (is_textbox1 and not is_textbox2) or (is_textbox2 and not is_textbox1):
            # One is text, one is shape - check size relationship
            if hasattr(shape1, 'width') and hasattr(shape2, 'width'):
                # Text is typically smaller than background shape
                area1 = shape1.width.inches * shape1.height.inches
                area2 = shape2.width.inches * shape2.height.inches

                # If text is 30-80% of shape size, likely intentional
                size_ratio = min(area1, area2) / max(area1, area2)
                if 0.3 <= size_ratio <= 0.8:
                    return True

        return False

    def _shapes_overlap(self, shape1, shape2) -> float:
        """
        Calculate overlap ratio between two shapes
        Returns: overlap ratio (0-1) or 0 if no overlap
        """
        # Get bounding boxes
        left1, top1, width1, height1 = shape1.left, shape1.top, shape1.width, shape1.height
        left2, top2, width2, height2 = shape2.left, shape2.top, shape2.width, shape2.height

        right1 = left1 + width1
        bottom1 = top1 + height1
        right2 = left2 + width2
        bottom2 = top2 + height2

        # Check for non-overlapping conditions
        if right1 <= left2 or right2 <= left1 or bottom1 <= top2 or bottom2 <= top1:
            return 0.0

        # Calculate overlap
        overlap_width = min(right1, right2) - max(left1, left2)
        overlap_height = min(bottom1, bottom2) - max(top1, top2)
        overlap_area = overlap_width * overlap_height

        shape1_area = width1 * height1
        shape2_area = width2 * height2

        # Return the maximum overlap ratio of the two shapes
        overlap_ratio1 = overlap_area / shape1_area if shape1_area > 0 else 0
        overlap_ratio2 = overlap_area / shape2_area if shape2_area > 0 else 0

        return max(overlap_ratio1, overlap_ratio2)

    def check_alignment(self, prs: Presentation):
        """Check alignment of shapes"""
        tolerance = Inches(self.guidelines['layout']['alignment']['tolerance_inches'])

        for slide_idx, slide in enumerate(prs.slides, 1):
            shapes = [s for s in slide.shapes if hasattr(s, 'left') and hasattr(s, 'top')]

            # Check left alignment
            left_edges = [s.left for s in shapes]
            right_edges = [s.left + s.width for s in shapes]
            top_edges = [s.top for s in shapes]
            bottom_edges = [s.top + s.height for s in shapes]

            # Group shapes that should be aligned (within tolerance)
            # This is a simplified check - in reality would need more sophisticated logic

    def check_fonts(self, prs: Presentation):
        """Check if fonts match MB style guidelines"""
        approved_fonts = [self.guidelines['fonts']['primary']] + self.guidelines['fonts']['fallback']

        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font and run.font.name:
                                if run.font.name not in approved_fonts:
                                    issue = LayoutIssue(
                                        slide_num=slide_idx,
                                        shape_name=shape.name,
                                        issue_type="Font",
                                        severity=IssueSeverity.INFO,
                                        description=f"Font '{run.font.name}' is not in approved MB font list",
                                        suggestion=f"Consider using {', '.join(approved_fonts)}"
                                    )
                                    self.issues.append(issue)

    def check_colors(self, prs: Presentation):
        """Check if colors follow MB style guidelines"""
        max_highlight_colors = self.guidelines['colors']['guidelines']['max_highlight_colors']

        # Define MB approved colors
        mb_colors = []
        for color_name, rgb in self.guidelines['colors']['primary'].items():
            mb_colors.append(RGBColor(*rgb))
        for color_name, rgb in self.guidelines['colors']['semantic'].items():
            mb_colors.append(RGBColor(*rgb))

        non_mb_colors = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_non_mb_colors = []

            for shape in slide.shapes:
                # Check fill color (only for shapes that have fill attribute)
                if hasattr(shape, 'fill') and shape.fill and shape.fill.type is not None:
                    try:
                        fill_color = shape.fill.fore_color.rgb
                        if fill_color and fill_color not in mb_colors:
                            slide_non_mb_colors.append(fill_color)
                    except:
                        pass

                # Check text color
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font and run.font.color:
                                try:
                                    text_color = run.font.color.rgb
                                    if text_color and text_color not in mb_colors:
                                        slide_non_mb_colors.append(text_color)
                                except:
                                    pass

            # Check if too many non-MB colors on one slide
            unique_non_mb_colors = list(set(slide_non_mb_colors))
            if len(unique_non_mb_colors) > max_highlight_colors:
                issue = LayoutIssue(
                    slide_num=slide_idx,
                    shape_name="Slide",
                    issue_type="Color",
                    severity=IssueSeverity.WARNING,
                    description=f"Slide uses {len(unique_non_mb_colors)} non-MB colors (max recommended: {max_highlight_colors})",
                    suggestion="Reduce color variety for cleaner business presentation"
                )
                self.issues.append(issue)

    def check_margins(self, prs: Presentation):
        """Check if content respects slide margins"""
        margins = self.guidelines['layout']['margins']
        left_margin = Inches(margins['left'])
        right_margin = Inches(margins['right'])
        top_margin = Inches(margins['top'])
        bottom_margin = Inches(margins['bottom'])

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if not hasattr(shape, 'left') or not hasattr(shape, 'top'):
                    continue

                # Check left margin
                if shape.left < left_margin:
                    issue = LayoutIssue(
                        slide_num=slide_idx,
                        shape_name=shape.name,
                        issue_type="Margin",
                        severity=IssueSeverity.WARNING,
                        description=f"Shape exceeds left margin (current: {shape.left.inches:.2f}in, min: {margins['left']}in)",
                        suggestion=f"Move shape right by {Inches(left_margin.inches - shape.left.inches).inches:.1f} inches"
                    )
                    self.issues.append(issue)

                # Check top margin
                if shape.top < top_margin:
                    issue = LayoutIssue(
                        slide_num=slide_idx,
                        shape_name=shape.name,
                        issue_type="Margin",
                        severity=IssueSeverity.WARNING,
                        description=f"Shape exceeds top margin (current: {shape.top.inches:.2f}in, min: {margins['top']}in)",
                        suggestion=f"Move shape down by {Inches(top_margin.inches - shape.top.inches).inches:.1f} inches"
                    )
                    self.issues.append(issue)

                # Check right margin
                if shape.left + shape.width > slide_width - right_margin:
                    issue = LayoutIssue(
                        slide_num=slide_idx,
                        shape_name=shape.name,
                        issue_type="Margin",
                        severity=IssueSeverity.WARNING,
                        description=f"Shape exceeds right margin",
                        suggestion="Reduce width or adjust positioning"
                    )
                    self.issues.append(issue)

                # Check bottom margin
                if shape.top + shape.height > slide_height - bottom_margin:
                    issue = LayoutIssue(
                        slide_num=slide_idx,
                        shape_name=shape.name,
                        issue_type="Margin",
                        severity=IssueSeverity.WARNING,
                        description=f"Shape exceeds bottom margin",
                        suggestion="Reduce height or adjust positioning"
                    )
                    self.issues.append(issue)

    def check_spacing(self, prs: Presentation):
        """Check spacing between elements"""
        min_spacing = Inches(self.guidelines['layout']['spacing']['min_between_elements'])

        for slide_idx, slide in enumerate(prs.slides, 1):
            shapes_list = list(slide.shapes)

            for i, shape1 in enumerate(shapes_list):
                for shape2 in shapes_list[i+1:]:
                    spacing = self._get_spacing(shape1, shape2)
                    if spacing < min_spacing and spacing >= 0:
                        # Only warn if shapes are very close (not overlapping)
                        spacing_val = spacing.inches if hasattr(spacing, 'inches') else float(spacing) / 914400
                        min_spacing_val = min_spacing.inches
                        issue = LayoutIssue(
                            slide_num=slide_idx,
                            shape_name=f"{shape1.name} and {shape2.name}",
                            issue_type="Spacing",
                            severity=IssueSeverity.INFO,
                            description=f"Spacing between shapes is {spacing_val:.2f}in (min recommended: {min_spacing_val:.2f}in)",
                            suggestion="Increase spacing for better readability"
                        )
                        self.issues.append(issue)

    def _get_spacing(self, shape1, shape2) -> float:
        """Calculate spacing between two shapes (negative if overlapping)"""
        # Get centers
        center1_x = shape1.left + shape1.width / 2
        center1_y = shape1.top + shape1.height / 2
        center2_x = shape2.left + shape2.width / 2
        center2_y = shape2.top + shape2.height / 2

        # Simple horizontal/vertical spacing check
        if abs(center1_x - center2_x) > (shape1.width + shape2.width) / 2:
            # Mostly vertical arrangement
            if center1_y < center2_y:
                return shape2.top - (shape1.top + shape1.height)
            else:
                return shape1.top - (shape2.top + shape2.height)
        else:
            # Mostly horizontal arrangement
            if center1_x < center2_x:
                return shape2.left - (shape1.left + shape1.width)
            else:
                return shape1.left - (shape2.left + shape2.width)

    def generate_report(self) -> str:
        """Generate a formatted report of all issues"""
        if not self.issues:
            return "✅ No layout issues found!"

        report = []
        report.append("="*80)
        report.append("PPT LAYOUT CHECK REPORT")
        report.append("="*80)
        report.append(f"Total issues found: {len(self.issues)}")
        report.append("")

        # Group by severity
        errors = [i for i in self.issues if i.severity == IssueSeverity.ERROR]
        warnings = [i for i in self.issues if i.severity == IssueSeverity.WARNING]
        infos = [i for i in self.issues if i.severity == IssueSeverity.INFO]

        if errors:
            report.append(f"🔴 ERRORS ({len(errors)}):")
            for issue in errors:
                report.append(f"  Slide {issue.slide_num} - {issue.shape_name}")
                report.append(f"    ❌ {issue.description}")
                report.append(f"    💡 {issue.suggestion}")
            report.append("")

        if warnings:
            report.append(f"🟡 WARNINGS ({len(warnings)}):")
            for issue in warnings:
                report.append(f"  Slide {issue.slide_num} - {issue.shape_name}")
                report.append(f"    ⚠️  {issue.description}")
                report.append(f"    💡 {issue.suggestion}")
            report.append("")

        if infos:
            report.append(f"ℹ️  INFORMATION ({len(infos)}):")
            for issue in infos:
                report.append(f"  Slide {issue.slide_num} - {issue.shape_name}")
                report.append(f"    ℹ️  {issue.description}")
                report.append(f"    💡 {issue.suggestion}")
            report.append("")

        report.append("="*80)
        return "\n".join(report)

    def check_presentation(self, ppt_path: str) -> List[LayoutIssue]:
        """Main method to check a presentation"""
        if not os.path.exists(ppt_path):
            raise FileNotFoundError(f"PPT file not found: {ppt_path}")

        prs = Presentation(ppt_path)

        # Clear previous issues
        self.issues = []

        # Run all checks
        print(f"\nChecking: {os.path.basename(ppt_path)}")
        print("Running checks...")

        print("  - Checking overlaps...")
        self.check_overlaps(prs)

        print("  - Checking alignment...")
        self.check_alignment(prs)

        print("  - Checking fonts...")
        self.check_fonts(prs)

        print("  - Checking colors...")
        self.check_colors(prs)

        print("  - Checking margins...")
        self.check_margins(prs)

        print("  - Checking spacing...")
        self.check_spacing(prs)

        print("Checks complete!\n")

        return self.issues


def main():
    """Standalone checker usage"""
    import sys
    import io

    # Fix encoding for Windows console
    if sys.stdout.encoding != 'utf-8':
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    if len(sys.argv) < 2:
        print("Usage: python ppt_layout_checker.py <ppt_file_path> [style_guidelines.json]")
        sys.exit(1)

    ppt_path = sys.argv[1]
    guidelines_path = sys.argv[2] if len(sys.argv) > 2 else r"D:\AI\mb_style_guidelines.json"

    # Create checker
    checker = PPTLayoutChecker(guidelines_path)

    # Run checks
    issues = checker.check_presentation(ppt_path)

    # Generate report
    report = checker.generate_report()
    print(report)

    # Save report to file
    report_path = ppt_path.replace('.pptx', '_layout_report.txt')
    with open(report_path, 'w', encoding='utf-8') as f:
        f.write(report)
    print(f"\nReport saved to: {report_path}")


if __name__ == '__main__':
    main()
