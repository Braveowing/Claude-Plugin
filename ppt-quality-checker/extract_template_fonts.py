"""
Extract font information from MB template files
"""
from pptx import Presentation
from pptx.util import Pt
import json

def extract_fonts_from_template(template_path):
    """Extract font information from a PowerPoint template"""
    prs = Presentation(template_path)
    font_info = {
        'title_fonts': [],
        'body_fonts': [],
        'all_fonts': set()
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_name = run.font.name
                        font_size = run.font.size.pt if run.font.size else None
                        font_bold = run.font.bold
                        font_italic = run.font.italic

                        font_data = {
                            'name': font_name,
                            'size': font_size,
                            'bold': font_bold,
                            'italic': font_italic
                        }

                        if font_name:
                            font_info['all_fonts'].add(font_name)

                            # Classify based on size
                            if font_size and font_size >= 24:
                                font_info['title_fonts'].append(font_data)
                            else:
                                font_info['body_fonts'].append(font_data)

    # Convert set to list for JSON serialization
    font_info['all_fonts'] = list(font_info['all_fonts'])

    return font_info

if __name__ == '__main__':
    # Extract from both template files
    templates = [
        r"D:\AI\参考模板\PPT template white.pptx",
        r"D:\AI\参考模板\Full template.pptx"
    ]

    all_font_info = {}

    for template_path in templates:
        try:
            print(f"\nExtracting fonts from: {template_path}")
            font_info = extract_fonts_from_template(template_path)

            print(f"  Fonts found: {font_info['all_fonts']}")
            print(f"  Title fonts: {len(font_info['title_fonts'])} instances")
            print(f"  Body fonts: {len(font_info['body_fonts'])} instances")

            # Save to JSON
            template_name = template_path.split('\\')[-1].replace(' ', '_').replace('.pptx', '')
            all_font_info[template_name] = font_info

        except Exception as e:
            print(f"  Error: {e}")

    # Save all font info
    output_path = r"D:\AI\mb_template_fonts.json"
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(all_font_info, f, indent=2, ensure_ascii=False)

    print(f"\nFont information saved to: {output_path}")
